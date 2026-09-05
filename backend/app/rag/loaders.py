import os
from pathlib import Path

from pypdf import PdfReader
from docx import Document
from pptx import Presentation


def load_pdf(file_path: str, document_id: int | None = None) -> str:
    reader = PdfReader(file_path)

    pages = []
    
    # Ensure static directory exists
    static_dir = None
    if document_id is not None:
        static_dir = Path(f"static/documents/{document_id}")
        static_dir.mkdir(parents=True, exist_ok=True)

    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        
        # Extract images if document_id is provided
        if document_id is not None and static_dir is not None:
            for image_file_object in page.images:
                img_name = image_file_object.name
                img_path = static_dir / img_name
                try:
                    with open(img_path, "wb") as fp:
                        fp.write(image_file_object.data)
                    # Inject marker into text
                    text += f"\n[PDF Diagram available at /static/documents/{document_id}/{img_name}]\n"
                except Exception as e:
                    print(f"Failed to save image {img_name}: {e}")

        pages.append(text)

    return "\n\n".join(pages)


def load_docx(file_path: str) -> str:
    document = Document(file_path)

    paragraphs = []

    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            paragraphs.append(paragraph.text.strip())

    return "\n\n".join(paragraphs)


def load_pptx(file_path: str) -> str:
    presentation = Presentation(file_path)

    slides = []

    for slide in presentation.slides:
        texts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

        if texts:
            slides.append("\n".join(texts))

    return "\n\n".join(slides)


def load_document(file_path: str, document_id: int | None = None) -> str:
    extension = Path(file_path).suffix.lower()

    if extension == ".pdf":
        return load_pdf(file_path, document_id)

    if extension == ".docx":
        return load_docx(file_path)

    if extension == ".pptx":
        return load_pptx(file_path)

    raise ValueError(f"Unsupported file type: {extension}")